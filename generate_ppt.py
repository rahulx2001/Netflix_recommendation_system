"""
Netflix Content Clustering - PowerPoint Presentation Generator
Generates a professional Netflix-themed PPT with real data from the analysis.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# CONFIG & CONSTANTS
# ============================================================
BASE = os.path.dirname(os.path.abspath(__file__))
OUT_FILE = os.path.join(BASE, "Netflix_Clustering_Presentation.pptx")

# Netflix brand colors
NETFLIX_RED = RGBColor(0xE5, 0x09, 0x14)
NETFLIX_BLACK = RGBColor(0x14, 0x14, 0x14)
NETFLIX_DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_YELLOW = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

# Chart image paths
CHARTS = {i: os.path.join(BASE, f) for i, f in enumerate([
    "01_missing_values.png",
    "02_content_type_distribution.png",
    "03_rating_distribution.png",
    "04_release_year_trend.png",
    "05_top_countries.png",
    "06_top_genres.png",
    "07_top_tfidf_terms.png",
    "08_pca_variance.png",
    "09_kmeans_metrics.png",
    "10_dendrogram.png",
    "11_cluster_visualization_2d.png",
    "12_cluster_sizes.png",
    "13_word_clouds.png",
    "14_cluster_composition.png",
], 1)}


# ============================================================
# HELPERS
# ============================================================
def set_slide_bg(slide, color=NETFLIX_BLACK):
    """Set the background fill color of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_red_bar(slide, top=Inches(0), height=Inches(0.06)):
    """Netflix-style thin red accent bar at top."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NETFLIX_RED
    shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Add a formatted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=16, color=LIGHT_GRAY, bold=False,
                  space_before=Pt(6), alignment=PP_ALIGN.LEFT):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add an image if the file exists."""
    if os.path.exists(path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(path, **kwargs)
    else:
        add_text_box(slide, left, top, Inches(4), Inches(0.5),
                     f"[Image not found: {os.path.basename(path)}]", 12, MUTED_GRAY)


def add_stat_box(slide, left, top, width, height, number, label, num_color=NETFLIX_RED):
    """Add a stat card with a big number and label below."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x35)
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x50)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = num_color
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)


def slide_number_footer(slide, num, total):
    """Add a small slide number at bottom right."""
    add_text_box(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num}/{total}", 10, MUTED_GRAY, alignment=PP_ALIGN.RIGHT)


def section_header_label(slide, text):
    """Small section label at top-left."""
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.35),
                 text, 11, NETFLIX_RED, bold=True)


# ============================================================
# BUILD PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

TOTAL_SLIDES = 15
blank_layout = prs.slide_layouts[6]  # blank

# =============================================================
# SLIDE 1 — TITLE SLIDE
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)

add_text_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
             "MACHINE LEARNING CASE STUDY", 14, NETFLIX_RED, bold=True)

tf = add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(2),
                  "Netflix Movies & TV Shows", 48, WHITE, bold=True)
add_paragraph(tf, "Content Clustering Analysis", 40, ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(1), Inches(4.2), Inches(10), Inches(0.8),
             "Unsupervised Machine Learning  •  NLP  •  TF-IDF  •  K-Means  •  Hierarchical Clustering",
             16, LIGHT_GRAY)

# Author & date
tf2 = add_text_box(slide, Inches(1), Inches(5.5), Inches(6), Inches(1.2),
                   "Rahul Kumar Singh", 22, WHITE, bold=True)
add_paragraph(tf2, "February 2026", 14, MUTED_GRAY)

# Accent stat boxes
add_stat_box(slide, Inches(8.5), Inches(5.2), Inches(1.8), Inches(1.1), "7,787", "Titles Analyzed")
add_stat_box(slide, Inches(10.6), Inches(5.2), Inches(1.8), Inches(1.1), "14", "Clusters Found", ACCENT_BLUE)

slide_number_footer(slide, 1, TOTAL_SLIDES)

# =============================================================
# SLIDE 2 — AGENDA / TABLE OF CONTENTS
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "OVERVIEW")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Agenda", 36, WHITE, bold=True)

agenda_items = [
    ("01", "Problem Statement", "Why content discovery needs ML"),
    ("02", "Dataset Overview", "7,787 Netflix titles with 12 features"),
    ("03", "Exploratory Data Analysis", "Missing values, distributions, trends"),
    ("04", "Text Preprocessing & NLP", "Tokenization, lemmatization, TF-IDF"),
    ("05", "Dimensionality Reduction", "PCA: 5,000 → 500 components"),
    ("06", "K-Means Clustering", "Optimal K=14 via Silhouette Score"),
    ("07", "Hierarchical Clustering", "Agglomerative validation with Ward linkage"),
    ("08", "Model Comparison", "K-Means vs Agglomerative on 3 metrics"),
    ("09", "Cluster Analysis", "14 cluster profiles with word clouds"),
    ("10", "Recommendation System", "Cosine similarity-based engine"),
    ("11", "Key Findings & Conclusions", "Results and business impact"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.7) + Inches(0.48) * i
    add_text_box(slide, Inches(1), y, Inches(0.6), Inches(0.4),
                 num, 14, NETFLIX_RED, bold=True)
    add_text_box(slide, Inches(1.7), y, Inches(4), Inches(0.4),
                 title, 16, WHITE, bold=True)
    add_text_box(slide, Inches(6), y, Inches(6), Inches(0.4),
                 desc, 14, MUTED_GRAY)

slide_number_footer(slide, 2, TOTAL_SLIDES)

# =============================================================
# SLIDE 3 — PROBLEM STATEMENT
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "01 — PROBLEM STATEMENT")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "The Challenge", 36, WHITE, bold=True)

tf = add_text_box(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(4),
                  "The Core Problem", 22, ACCENT_BLUE, bold=True)
add_paragraph(tf, "• Netflix has 7,787+ titles across diverse genres, languages, and formats", 15, LIGHT_GRAY)
add_paragraph(tf, "• Users face choice paralysis — too many options, not enough structure", 15, LIGHT_GRAY)
add_paragraph(tf, "• Traditional recommenders rely on viewing history — fails for new users", 15, LIGHT_GRAY)
add_paragraph(tf, "• Cold-start problem: newly added content has no interaction data", 15, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)

p = add_paragraph(tf, "Our Solution", 22, ACCENT_GREEN, bold=True, space_before=Pt(18))
add_paragraph(tf, "Content-based clustering using NLP + unsupervised ML", 15, LIGHT_GRAY)
add_paragraph(tf, "→ Group similar content by text features (description, cast, genre, director)", 15, LIGHT_GRAY)
add_paragraph(tf, "→ Recommend within clusters using cosine similarity", 15, LIGHT_GRAY)
add_paragraph(tf, "→ No user data needed — works from day one", 15, LIGHT_GRAY)

# Pipeline box
tf2 = add_text_box(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(1),
                   "ML Pipeline", 18, WHITE, bold=True)
add_paragraph(tf2, "Raw Data → Cleaning → NLP → TF-IDF → PCA → Clustering → Recommendations", 14, ACCENT_YELLOW, bold=True)

# Stat boxes
add_stat_box(slide, Inches(7), Inches(3.3), Inches(2.5), Inches(1.2), "5,000", "TF-IDF Features", ACCENT_BLUE)
add_stat_box(slide, Inches(10), Inches(3.3), Inches(2.5), Inches(1.2), "500", "PCA Components", ACCENT_GREEN)
add_stat_box(slide, Inches(7), Inches(4.8), Inches(2.5), Inches(1.2), "14", "Clusters Found", NETFLIX_RED)
add_stat_box(slide, Inches(10), Inches(4.8), Inches(2.5), Inches(1.2), "7,787", "Titles Analyzed", ACCENT_YELLOW)

slide_number_footer(slide, 3, TOTAL_SLIDES)

# =============================================================
# SLIDE 4 — DATASET OVERVIEW
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "02 — DATASET OVERVIEW")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Netflix Content Dataset", 36, WHITE, bold=True)

tf = add_text_box(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(5),
                  "Dataset Characteristics", 20, ACCENT_BLUE, bold=True)
add_paragraph(tf, "• 7,787 titles (Movies & TV Shows)", 15, LIGHT_GRAY)
add_paragraph(tf, "• 12 original features", 15, LIGHT_GRAY)
add_paragraph(tf, "• Collected up to 2021", 15, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "Key Features Used for Clustering:", 16, WHITE, bold=True, space_before=Pt(12))
add_paragraph(tf, "  title  •  director  •  cast  •  country", 14, LIGHT_GRAY)
add_paragraph(tf, "  listed_in (genres)  •  description", 14, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "Content Breakdown:", 16, WHITE, bold=True, space_before=Pt(12))
add_paragraph(tf, "  Movies:    5,377 (69%)", 14, LIGHT_GRAY)
add_paragraph(tf, "  TV Shows:  2,410 (31%)", 14, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "Missing Values (Top 3):", 16, WHITE, bold=True, space_before=Pt(12))
add_paragraph(tf, "  director — 2,389 (30.7%)", 14, LIGHT_GRAY)
add_paragraph(tf, "  cast — 718 (9.2%)", 14, LIGHT_GRAY)
add_paragraph(tf, "  country — 507 (6.5%)", 14, LIGHT_GRAY)

add_image_safe(slide, CHARTS[1], Inches(7), Inches(1.8), width=Inches(5.5))
slide_number_footer(slide, 4, TOTAL_SLIDES)

# =============================================================
# SLIDE 5 — EDA: DISTRIBUTIONS
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "03 — EXPLORATORY DATA ANALYSIS")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Content Distribution Analysis", 36, WHITE, bold=True)

add_image_safe(slide, CHARTS[2], Inches(0.5), Inches(1.8), width=Inches(6))
add_image_safe(slide, CHARTS[3], Inches(6.8), Inches(1.8), width=Inches(6))

tf = add_text_box(slide, Inches(0.5), Inches(5.8), Inches(6), Inches(1.5),
                  "Content Type", 14, ACCENT_BLUE, bold=True)
add_paragraph(tf, "69% Movies vs 31% TV Shows — movies dominate the catalog", 13, LIGHT_GRAY)

tf2 = add_text_box(slide, Inches(6.8), Inches(5.8), Inches(6), Inches(1.5),
                   "Rating Distribution", 14, ACCENT_BLUE, bold=True)
add_paragraph(tf2, "TV-MA is the most common rating — Netflix focuses on adult content", 13, LIGHT_GRAY)

slide_number_footer(slide, 5, TOTAL_SLIDES)

# =============================================================
# SLIDE 6 — EDA: TRENDS & GEOGRAPHY
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "03 — EXPLORATORY DATA ANALYSIS")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Content Trends & Geography", 36, WHITE, bold=True)

add_image_safe(slide, CHARTS[4], Inches(0.5), Inches(1.8), width=Inches(6))
add_image_safe(slide, CHARTS[5], Inches(6.8), Inches(1.8), width=Inches(6))

tf = add_text_box(slide, Inches(0.5), Inches(5.8), Inches(6), Inches(1.5),
                  "Release Year Trend", 14, ACCENT_BLUE, bold=True)
add_paragraph(tf, "10x growth from 2015-2020 driven by Netflix Originals expansion", 13, LIGHT_GRAY)

tf2 = add_text_box(slide, Inches(6.8), Inches(5.8), Inches(6), Inches(1.5),
                   "Top Countries", 14, ACCENT_BLUE, bold=True)
add_paragraph(tf2, "US leads production, but India & UK contribute significantly — 35% international", 13, LIGHT_GRAY)

slide_number_footer(slide, 6, TOTAL_SLIDES)

# =============================================================
# SLIDE 7 — EDA: GENRES
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "03 — EXPLORATORY DATA ANALYSIS")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Genre Distribution", 36, WHITE, bold=True)

add_image_safe(slide, CHARTS[6], Inches(1.5), Inches(1.8), width=Inches(5.5))

tf = add_text_box(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(5),
                  "Key Observations", 20, ACCENT_BLUE, bold=True)
add_paragraph(tf, "• International Movies is the dominant genre", 15, LIGHT_GRAY)
add_paragraph(tf, "• Dramas appear in 50%+ of all titles", 15, LIGHT_GRAY)
add_paragraph(tf, "• Comedies and Documentaries follow closely", 15, LIGHT_GRAY)
add_paragraph(tf, "• Multi-genre labels are common (e.g., 'International Movies, Dramas, Comedies')", 15, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "Implication for Clustering:", 16, WHITE, bold=True, space_before=Pt(12))
add_paragraph(tf, "Genre alone isn't enough — we combine description, cast, director, and country to capture deeper content patterns", 14, LIGHT_GRAY)

slide_number_footer(slide, 7, TOTAL_SLIDES)

# =============================================================
# SLIDE 8 — NLP & TF-IDF
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "04 — TEXT PREPROCESSING & TF-IDF")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "NLP Pipeline & TF-IDF Vectorization", 36, WHITE, bold=True)

tf = add_text_box(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(5),
                  "Text Processing Steps", 20, ACCENT_BLUE, bold=True)
add_paragraph(tf, "1. Combined 6 text columns into one:", 15, LIGHT_GRAY)
add_paragraph(tf, "   title + director + cast + country + genre + description", 13, MUTED_GRAY)
add_paragraph(tf, "2. Lowercased, removed punctuation & numbers", 15, LIGHT_GRAY)
add_paragraph(tf, "3. Tokenization using NLTK word_tokenize", 15, LIGHT_GRAY)
add_paragraph(tf, "4. Removed stopwords (English + custom set)", 15, LIGHT_GRAY)
add_paragraph(tf, "5. Lemmatization with WordNetLemmatizer", 15, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "TF-IDF Configuration", 18, ACCENT_GREEN, bold=True, space_before=Pt(14))
add_paragraph(tf, "• max_features = 5,000", 15, LIGHT_GRAY)
add_paragraph(tf, "• ngram_range = (1, 2) — unigrams + bigrams", 15, LIGHT_GRAY)
add_paragraph(tf, "• min_df = 5, max_df = 0.8", 15, LIGHT_GRAY)
add_paragraph(tf, "• Result: 7,787 × 5,000 sparse matrix", 15, LIGHT_GRAY)

add_image_safe(slide, CHARTS[7], Inches(6.8), Inches(1.8), width=Inches(6))
slide_number_footer(slide, 8, TOTAL_SLIDES)

# =============================================================
# SLIDE 9 — PCA
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "05 — DIMENSIONALITY REDUCTION")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "PCA Dimensionality Reduction", 36, WHITE, bold=True)

add_image_safe(slide, CHARTS[8], Inches(0.5), Inches(1.8), width=Inches(6.5))

tf = add_text_box(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(5),
                  "Why PCA?", 20, ACCENT_BLUE, bold=True)
add_paragraph(tf, "• 5,000 TF-IDF features is too high-dimensional", 15, LIGHT_GRAY)
add_paragraph(tf, "• Curse of dimensionality makes distances meaningless", 15, LIGHT_GRAY)
add_paragraph(tf, "• Clustering algorithms struggle in high-D space", 15, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "PCA Results", 18, ACCENT_GREEN, bold=True, space_before=Pt(14))
add_paragraph(tf, "• 95% variance needs 3,111 components", 15, LIGHT_GRAY)
add_paragraph(tf, "• Capped at 500 components (practical limit)", 15, LIGHT_GRAY)
add_paragraph(tf, "• 42.77% variance explained with 500 components", 15, LIGHT_GRAY)
add_paragraph(tf, "• Reduced shape: 7,787 × 500", 15, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "Note: Text data is inherently high-dimensional — even 500 components capture substantial structure for clustering.", 13, MUTED_GRAY)

slide_number_footer(slide, 9, TOTAL_SLIDES)

# =============================================================
# SLIDE 10 — K-MEANS CLUSTERING
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "06 — K-MEANS CLUSTERING")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Finding Optimal K", 36, WHITE, bold=True)

add_image_safe(slide, CHARTS[9], Inches(0.3), Inches(1.7), width=Inches(7))

tf = add_text_box(slide, Inches(7.5), Inches(1.8), Inches(5.3), Inches(5),
                  "Evaluation Metrics (K=2 to 14)", 18, ACCENT_BLUE, bold=True)
add_paragraph(tf, "• Elbow Method — Inertia decreases steadily, no sharp elbow", 14, LIGHT_GRAY)
add_paragraph(tf, "• Silhouette Score — K=14 achieves highest score (0.0338)", 14, LIGHT_GRAY)
add_paragraph(tf, "• Calinski-Harabasz — Higher is better (73.65 at K=14)", 14, LIGHT_GRAY)
add_paragraph(tf, "• Davies-Bouldin — Lower is better (4.80 at K=14)", 14, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "Final K-Means Configuration", 18, ACCENT_GREEN, bold=True, space_before=Pt(14))
add_paragraph(tf, f"• n_clusters = 14", 14, LIGHT_GRAY)
add_paragraph(tf, f"• n_init = 20, max_iter = 500", 14, LIGHT_GRAY)
add_paragraph(tf, f"• random_state = 42", 14, LIGHT_GRAY)

add_stat_box(slide, Inches(7.5), Inches(5.5), Inches(2.4), Inches(1.1), "0.0338", "Silhouette Score", NETFLIX_RED)
add_stat_box(slide, Inches(10.3), Inches(5.5), Inches(2.4), Inches(1.1), "K = 14", "Optimal Clusters", ACCENT_BLUE)

slide_number_footer(slide, 10, TOTAL_SLIDES)

# =============================================================
# SLIDE 11 — HIERARCHICAL CLUSTERING & COMPARISON
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "07 & 08 — HIERARCHICAL CLUSTERING & COMPARISON")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Agglomerative Clustering & Model Comparison", 32, WHITE, bold=True)

add_image_safe(slide, CHARTS[10], Inches(0.3), Inches(1.8), width=Inches(6.5))

tf = add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(1),
                  "Comparison: K-Means vs Agglomerative", 18, ACCENT_BLUE, bold=True)

# Comparison table header
y_start = Inches(2.8)
row_h = Inches(0.45)
cols = [Inches(7.2), Inches(9.5), Inches(11.2)]
col_widths = [Inches(2.3), Inches(1.7), Inches(1.7)]

headers = ["Metric", "K-Means", "Agglomerative"]
for j, (hdr, cx, cw) in enumerate(zip(headers, cols, col_widths)):
    add_text_box(slide, cx, y_start, cw, row_h, hdr, 13, NETFLIX_RED, bold=True)

rows = [
    ("Silhouette ↑", "0.0338", "0.0186"),
    ("Calinski-Harabasz ↑", "73.65", "61.01"),
    ("Davies-Bouldin ↓", "4.80", "4.94"),
]
for i, (metric, km, ag) in enumerate(rows):
    y = y_start + row_h * (i + 1)
    add_text_box(slide, cols[0], y, col_widths[0], row_h, metric, 13, LIGHT_GRAY)
    add_text_box(slide, cols[1], y, col_widths[1], row_h, km, 13, ACCENT_GREEN, bold=True)
    add_text_box(slide, cols[2], y, col_widths[2], row_h, ag, 13, MUTED_GRAY)

tf2 = add_text_box(slide, Inches(7.2), Inches(4.8), Inches(5.5), Inches(2),
                   "🏆 Winner: K-Means", 22, ACCENT_GREEN, bold=True)
add_paragraph(tf2, "K-Means outperforms on all 3 evaluation metrics.", 14, LIGHT_GRAY)
add_paragraph(tf2, "Selected as the final clustering method.", 14, LIGHT_GRAY)
add_paragraph(tf2, "Ward linkage Agglomerative used as validation — confirms structure exists.", 14, MUTED_GRAY)

slide_number_footer(slide, 11, TOTAL_SLIDES)

# =============================================================
# SLIDE 12 — CLUSTER VISUALIZATION
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "09 — CLUSTER ANALYSIS")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "2D Cluster Visualization & Sizes", 36, WHITE, bold=True)

add_image_safe(slide, CHARTS[11], Inches(0.3), Inches(1.7), width=Inches(6.5))
add_image_safe(slide, CHARTS[12], Inches(7), Inches(1.7), width=Inches(5.8))

tf = add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(1.2),
                  "Observation:", 14, ACCENT_BLUE, bold=True)
add_paragraph(tf, "Cluster sizes range from 148 (niche Spanish-Language TV) to 1,433 (general Dramas) — natural distribution reflecting Netflix's diverse catalog.", 13, LIGHT_GRAY)

slide_number_footer(slide, 12, TOTAL_SLIDES)

# =============================================================
# SLIDE 13 — CLUSTER PROFILES
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "09 — CLUSTER PROFILES")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "14 Cluster Profiles", 36, WHITE, bold=True)

# Cluster data from cluster_summary.csv
clusters = [
    (0, 387, "Kids' TV", "16 Movies, 371 TV Shows"),
    (1, 770, "Documentaries", "768 Movies, 2 TV Shows"),
    (2, 837, "Children & Family", "794 Movies, 43 TV Shows"),
    (3, 225, "British TV Shows", "0 Movies, 225 TV Shows"),
    (4, 316, "Romantic Movies", "316 Movies, 0 TV Shows"),
    (5, 688, "Action & Adventure", "588 Movies, 100 TV Shows"),
    (6, 595, "International Movies", "595 Movies, 0 TV Shows"),
    (7, 1433, "Dramas", "1,186 Movies, 247 TV Shows"),
    (8, 685, "International Movies (B)", "683 Movies, 2 TV Shows"),
    (9, 397, "International Movies (C)", "397 Movies, 0 TV Shows"),
    (10, 339, "Docuseries", "0 Movies, 339 TV Shows"),
    (11, 787, "International TV Shows", "1 Movie, 786 TV Shows"),
    (12, 180, "Int'l TV Shows (B)", "33 Movies, 147 TV Shows"),
    (13, 148, "Spanish-Language TV", "0 Movies, 148 TV Shows"),
]

col1_x = Inches(0.5)
col2_x = Inches(6.8)
per_col = 7

for i, (cid, size, genre, breakdown) in enumerate(clusters):
    col = 0 if i < per_col else 1
    row = i if i < per_col else i - per_col
    x = col1_x if col == 0 else col2_x
    y = Inches(1.7) + Inches(0.72) * row

    # Cluster number badge
    colors = [NETFLIX_RED, ACCENT_BLUE, ACCENT_GREEN, ACCENT_YELLOW,
              RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xEC, 0x48, 0x99),
              RGBColor(0x06, 0xB6, 0xD4), NETFLIX_RED, ACCENT_BLUE, ACCENT_GREEN,
              ACCENT_YELLOW, RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xEC, 0x48, 0x99),
              RGBColor(0x06, 0xB6, 0xD4)]

    add_text_box(slide, x, y, Inches(0.5), Inches(0.35),
                 f"C{cid}", 11, colors[i], bold=True)
    tf = add_text_box(slide, x + Inches(0.5), y, Inches(5.5), Inches(0.35),
                      f"{genre}  ({size} titles)", 14, WHITE, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.3), Inches(5.5), Inches(0.3),
                 breakdown, 11, MUTED_GRAY)

slide_number_footer(slide, 13, TOTAL_SLIDES)

# =============================================================
# SLIDE 14 — WORD CLOUDS + COMPOSITION
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "09 — CLUSTER WORD CLOUDS")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Cluster Word Clouds & Composition", 36, WHITE, bold=True)

add_image_safe(slide, CHARTS[13], Inches(0.3), Inches(1.7), width=Inches(7))
add_image_safe(slide, CHARTS[14], Inches(7.5), Inches(1.7), width=Inches(5.3))

tf = add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(1.2),
                  "Validation:", 14, ACCENT_GREEN, bold=True)
add_paragraph(tf, "Distinct word clouds per cluster confirm the model captured meaningful thematic groupings — not random noise.", 13, LIGHT_GRAY)

slide_number_footer(slide, 14, TOTAL_SLIDES)

# =============================================================
# SLIDE 15 — CONCLUSIONS & KEY FINDINGS
# =============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_red_bar(slide)
section_header_label(slide, "11 — CONCLUSIONS")

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Key Findings & Conclusions", 36, WHITE, bold=True)

tf = add_text_box(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(5),
                  "What We Achieved", 22, ACCENT_BLUE, bold=True)
add_paragraph(tf, "✓  Clustered 7,787 Netflix titles into 14 meaningful groups", 15, LIGHT_GRAY)
add_paragraph(tf, "✓  K-Means outperformed Agglomerative on all 3 metrics", 15, LIGHT_GRAY)
add_paragraph(tf, "✓  Built a content-based recommendation engine using cosine similarity", 15, LIGHT_GRAY)
add_paragraph(tf, "✓  Solved the cold-start problem — no user history needed", 15, LIGHT_GRAY)
add_paragraph(tf, "✓  Distinct word clouds validate cluster quality", 15, LIGHT_GRAY)
add_paragraph(tf, "", 8, LIGHT_GRAY)
add_paragraph(tf, "Technical Summary", 20, ACCENT_GREEN, bold=True, space_before=Pt(14))
add_paragraph(tf, "• NLP: Tokenization + Lemmatization + Stopword removal", 14, LIGHT_GRAY)
add_paragraph(tf, "• TF-IDF: 5,000 features (unigrams + bigrams)", 14, LIGHT_GRAY)
add_paragraph(tf, "• PCA: 5,000 → 500 components (42.77% variance)", 14, LIGHT_GRAY)
add_paragraph(tf, "• K-Means: K=14, Silhouette=0.034, CH=73.65, DB=4.80", 14, LIGHT_GRAY)
add_paragraph(tf, "• Recommendation: Cosine similarity on TF-IDF vectors", 14, LIGHT_GRAY)

tf2 = add_text_box(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                   "Business Applications", 22, ACCENT_YELLOW, bold=True)
add_paragraph(tf2, "📺  Content Discovery", 16, WHITE, bold=True, space_before=Pt(10))
add_paragraph(tf2, "Group similar content for better browsing experience", 14, LIGHT_GRAY)
add_paragraph(tf2, "", 6, LIGHT_GRAY)
add_paragraph(tf2, "🆕  Cold-Start Recommendations", 16, WHITE, bold=True, space_before=Pt(10))
add_paragraph(tf2, "Recommend new/low-traffic content via cluster membership", 14, LIGHT_GRAY)
add_paragraph(tf2, "", 6, LIGHT_GRAY)
add_paragraph(tf2, "📊  Content Strategy", 16, WHITE, bold=True, space_before=Pt(10))
add_paragraph(tf2, "Identify gaps and trends across content clusters", 14, LIGHT_GRAY)
add_paragraph(tf2, "", 6, LIGHT_GRAY)
add_paragraph(tf2, "🌍  International Content", 16, WHITE, bold=True, space_before=Pt(10))
add_paragraph(tf2, "Multiple international clusters suggest localized recommendation potential", 14, LIGHT_GRAY)

# Thank you line
add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
             "Thank You  •  Rahul Kumar Singh  •  February 2026", 14, MUTED_GRAY,
             alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 15, TOTAL_SLIDES)

# ============================================================
# SAVE
# ============================================================
prs.save(OUT_FILE)
print(f"✅ Presentation saved: {OUT_FILE}")
print(f"   {TOTAL_SLIDES} slides generated with Netflix theme")
